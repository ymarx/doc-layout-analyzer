"""
Content Extractor - 콘텐츠 추출 파이프라인
텍스트, 표, 다이어그램, 이미지 등 다양한 콘텐츠 추출
"""

import logging
import time
import asyncio
from pathlib import Path
from typing import Dict, List, Any, Optional, Union, Tuple
import numpy as np
from PIL import Image
import io
import json
from dataclasses import dataclass
from enum import Enum

# OCR 관련
try:
    from paddleocr import PaddleOCR
    PADDLEOCR_AVAILABLE = True
except ImportError:
    PADDLEOCR_AVAILABLE = False

# 표 추출 관련
try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError:
    CAMELOT_AVAILABLE = False

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

# PPTX 다이어그램 처리
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

from ..analyzers.layout_analyzer import LayoutElement, LayoutElementType
from ..core.device_manager import DeviceManager
from ..core.config import ConfigManager

logger = logging.getLogger(__name__)


@dataclass
class ExtractedContent:
    """추출된 콘텐츠"""
    content_type: str
    data: Any
    bbox: Optional[List[float]] = None
    page_number: int = 0
    confidence: float = 1.0
    metadata: Optional[Dict[str, Any]] = None


@dataclass
class ExtractionResult:
    """추출 결과"""
    success: bool
    content: List[ExtractedContent] = None
    processing_time: float = 0.0
    method_used: str = ""
    error: Optional[str] = None


class ContentExtractor:
    """콘텐츠 추출 파이프라인"""

    def __init__(self, device_manager: DeviceManager = None, config: ConfigManager = None):
        self.device_manager = device_manager
        self.config = config
        self.current_device = "cpu"

        # 추출기 초기화
        self.text_extractor = None
        self.table_extractor = None
        self.diagram_extractor = None

        self._initialize_extractors()

        logger.info("Content Extractor 초기화 완료")

    def _initialize_extractors(self):
        """추출기 초기화"""
        # 디바이스 설정
        use_gpu = False
        if self.device_manager and self.config:
            self.current_device = self.device_manager.get_optimal_device(prefer_gpu=True)
            use_gpu = self.current_device != "cpu" and self.config.ocr.use_gpu

        # 텍스트 추출기 (PaddleOCR)
        if PADDLEOCR_AVAILABLE:
            try:
                self.text_extractor = PaddleOCR(
                    use_angle_cls=True,
                    lang="korean"
                )
                logger.info(f"PaddleOCR 텍스트 추출기 초기화 완료 (GPU: {use_gpu})")
            except Exception as e:
                logger.error(f"PaddleOCR 초기화 실패: {e}")

    async def extract_from_layout(self,
                                 layout_elements: List[LayoutElement],
                                 document_path: Union[str, Path],
                                 page_images: List[np.ndarray] = None) -> ExtractionResult:
        """레이아웃 분석 결과를 바탕으로 콘텐츠 추출"""
        start_time = time.time()

        try:
            extracted_content = []

            # 요소 타입별로 그룹화
            elements_by_type = self._group_elements_by_type(layout_elements)

            # 병렬 추출 작업
            tasks = []

            # 텍스트 추출
            if elements_by_type.get(LayoutElementType.TEXT):
                tasks.append(self._extract_text_elements(
                    elements_by_type[LayoutElementType.TEXT],
                    page_images
                ))

            # 표 추출
            if elements_by_type.get(LayoutElementType.TABLE):
                tasks.append(self._extract_table_elements(
                    elements_by_type[LayoutElementType.TABLE],
                    document_path
                ))

            # 다이어그램 추출
            if elements_by_type.get(LayoutElementType.DIAGRAM):
                tasks.append(self._extract_diagram_elements(
                    elements_by_type[LayoutElementType.DIAGRAM],
                    document_path
                ))

            # 이미지 추출
            if elements_by_type.get(LayoutElementType.IMAGE):
                tasks.append(self._extract_image_elements(
                    elements_by_type[LayoutElementType.IMAGE],
                    document_path
                ))

            # 병렬 실행
            if tasks:
                results = await asyncio.gather(*tasks, return_exceptions=True)

                # 결과 병합
                for result in results:
                    if isinstance(result, Exception):
                        logger.error(f"추출 작업 실패: {result}")
                        continue

                    if isinstance(result, list):
                        extracted_content.extend(result)

            processing_time = time.time() - start_time

            return ExtractionResult(
                success=True,
                content=extracted_content,
                processing_time=processing_time,
                method_used="multi_extractor"
            )

        except Exception as e:
            logger.error(f"콘텐츠 추출 실패: {e}")
            return ExtractionResult(
                success=False,
                error=str(e),
                processing_time=time.time() - start_time
            )

    def _group_elements_by_type(self, elements: List[LayoutElement]) -> Dict[LayoutElementType, List[LayoutElement]]:
        """레이아웃 요소를 타입별로 그룹화"""
        grouped = {}

        for element in elements:
            element_type = element.element_type
            if element_type not in grouped:
                grouped[element_type] = []
            grouped[element_type].append(element)

        return grouped

    async def _extract_text_elements(self,
                                   text_elements: List[LayoutElement],
                                   page_images: List[np.ndarray] = None) -> List[ExtractedContent]:
        """텍스트 요소 추출"""
        extracted_texts = []

        if not self.text_extractor:
            logger.warning("텍스트 추출기를 사용할 수 없습니다.")
            return extracted_texts

        try:
            for element in text_elements:
                # 이미 텍스트가 있는 경우
                if element.text:
                    content = ExtractedContent(
                        content_type="text",
                        data=element.text,
                        bbox=element.bbox,
                        page_number=element.page_number,
                        confidence=element.confidence,
                        metadata={
                            "extraction_method": "layout_based",
                            "original_element_type": element.element_type.value if hasattr(element.element_type, 'value') else str(element.element_type)
                        }
                    )
                    extracted_texts.append(content)
                    continue

                # OCR로 텍스트 추출 (페이지 이미지가 있는 경우)
                if page_images and element.page_number <= len(page_images):
                    page_image = page_images[element.page_number - 1]
                    text = await self._extract_text_from_region(page_image, element.bbox)

                    if text:
                        content = ExtractedContent(
                            content_type="text",
                            data=text,
                            bbox=element.bbox,
                            page_number=element.page_number,
                            confidence=element.confidence,
                            metadata={
                                "extraction_method": "ocr",
                                "original_element_type": element.element_type.value if hasattr(element.element_type, 'value') else str(element.element_type)
                            }
                        )
                        extracted_texts.append(content)

        except Exception as e:
            logger.error(f"텍스트 추출 실패: {e}")

        return extracted_texts

    async def _extract_text_from_region(self, page_image: np.ndarray, bbox: List[float]) -> Optional[str]:
        """이미지 영역에서 텍스트 추출"""
        try:
            # 영역 크롭
            x1, y1, x2, y2 = [int(coord) for coord in bbox]

            # 경계 확인
            h, w = page_image.shape[:2]
            x1 = max(0, min(x1, w))
            y1 = max(0, min(y1, h))
            x2 = max(x1, min(x2, w))
            y2 = max(y1, min(y2, h))

            cropped_image = page_image[y1:y2, x1:x2]

            if cropped_image.size == 0:
                return None

            # OCR 수행
            result = await asyncio.to_thread(self.text_extractor.ocr, cropped_image)

            # 결과 파싱
            text_parts = []
            for line in result:
                if isinstance(line, list) and len(line) >= 2:
                    text_info = line[1]
                    if isinstance(text_info, tuple) and len(text_info) >= 1:
                        text_parts.append(text_info[0])
                    else:
                        text_parts.append(str(text_info))

            return ' '.join(text_parts).strip() if text_parts else None

        except Exception as e:
            logger.debug(f"영역 텍스트 추출 실패: {e}")
            return None

    async def _extract_table_elements(self,
                                    table_elements: List[LayoutElement],
                                    document_path: Union[str, Path]) -> List[ExtractedContent]:
        """표 요소 추출"""
        extracted_tables = []

        try:
            document_path = Path(document_path)

            # PDF인 경우 Camelot과 pdfplumber 사용
            if document_path.suffix.lower() == '.pdf':
                extracted_tables = await self._extract_tables_from_pdf(table_elements, document_path)

            # 다른 형식의 경우 이미지 기반 추출
            else:
                # TODO: 이미지 기반 표 추출 구현
                pass

        except Exception as e:
            logger.error(f"표 추출 실패: {e}")

        return extracted_tables

    async def _extract_tables_from_pdf(self,
                                     table_elements: List[LayoutElement],
                                     pdf_path: Path) -> List[ExtractedContent]:
        """PDF에서 표 추출"""
        extracted_tables = []

        try:
            # 페이지별 표 요소 그룹화
            tables_by_page = {}
            for element in table_elements:
                page_num = element.page_number
                if page_num not in tables_by_page:
                    tables_by_page[page_num] = []
                tables_by_page[page_num].append(element)

            # 각 페이지별로 처리
            for page_num, page_tables in tables_by_page.items():
                # Camelot으로 표 추출
                camelot_tables = []
                if CAMELOT_AVAILABLE:
                    camelot_tables = await self._extract_tables_with_camelot(pdf_path, page_num, page_tables)

                # pdfplumber로 표 추출 (백업)
                pdfplumber_tables = []
                if PDFPLUMBER_AVAILABLE:
                    pdfplumber_tables = await self._extract_tables_with_pdfplumber(pdf_path, page_num, page_tables)

                # 결과 병합 및 검증
                best_tables = self._select_best_tables(camelot_tables, pdfplumber_tables, page_tables)
                extracted_tables.extend(best_tables)

        except Exception as e:
            logger.error(f"PDF 표 추출 실패: {e}")

        return extracted_tables

    async def _extract_tables_with_camelot(self,
                                         pdf_path: Path,
                                         page_num: int,
                                         table_elements: List[LayoutElement]) -> List[ExtractedContent]:
        """Camelot으로 표 추출"""
        tables = []

        try:
            # 표 영역들을 문자열로 변환 (Camelot 형식)
            table_areas = []
            for element in table_elements:
                x1, y1, x2, y2 = element.bbox
                # Camelot은 PDF 좌표계 사용 (하단 좌측 원점)
                area = f"{x1},{y1},{x2},{y2}"
                table_areas.append(area)

            if not table_areas:
                return tables

            # Camelot 실행 (먼저 lattice 모드)
            camelot_result = await asyncio.to_thread(
                camelot.read_pdf,
                str(pdf_path),
                pages=str(page_num),
                flavor='lattice',
                table_areas=table_areas
            )

            # 결과 처리
            for i, table in enumerate(camelot_result):
                if table.accuracy > 0.8:  # 높은 정확도만
                    content = ExtractedContent(
                        content_type="table",
                        data={
                            "data": table.df.values.tolist(),
                            "headers": table.df.columns.tolist(),
                            "shape": table.df.shape,
                            "html": table.df.to_html(),
                            "csv": table.df.to_csv()
                        },
                        bbox=table_elements[i].bbox if i < len(table_elements) else None,
                        page_number=page_num,
                        confidence=table.accuracy,
                        metadata={
                            "extraction_method": "camelot_lattice",
                            "accuracy": table.accuracy
                        }
                    )
                    tables.append(content)

            # 정확도가 낮으면 stream 모드 시도
            if not tables or all(t.confidence < 0.8 for t in tables):
                camelot_result = await asyncio.to_thread(
                    camelot.read_pdf,
                    str(pdf_path),
                    pages=str(page_num),
                    flavor='stream'
                )

                for i, table in enumerate(camelot_result):
                    if table.accuracy > 0.7:
                        content = ExtractedContent(
                            content_type="table",
                            data={
                                "data": table.df.values.tolist(),
                                "headers": table.df.columns.tolist(),
                                "shape": table.df.shape,
                                "html": table.df.to_html(),
                                "csv": table.df.to_csv()
                            },
                            bbox=table_elements[i].bbox if i < len(table_elements) else None,
                            page_number=page_num,
                            confidence=table.accuracy,
                            metadata={
                                "extraction_method": "camelot_stream",
                                "accuracy": table.accuracy
                            }
                        )
                        tables.append(content)

        except Exception as e:
            logger.debug(f"Camelot 표 추출 실패 (페이지 {page_num}): {e}")

        return tables

    async def _extract_tables_with_pdfplumber(self,
                                            pdf_path: Path,
                                            page_num: int,
                                            table_elements: List[LayoutElement]) -> List[ExtractedContent]:
        """pdfplumber로 표 추출"""
        tables = []

        try:
            with pdfplumber.open(str(pdf_path)) as pdf:
                if page_num <= len(pdf.pages):
                    page = pdf.pages[page_num - 1]  # 0-based index

                    # 페이지에서 모든 표 추출
                    page_tables = page.extract_tables()

                    for i, table_data in enumerate(page_tables):
                        if table_data and any(any(cell for cell in row if cell) for row in table_data):
                            # 해당하는 레이아웃 요소 찾기
                            corresponding_element = table_elements[i] if i < len(table_elements) else None

                            content = ExtractedContent(
                                content_type="table",
                                data={
                                    "data": table_data,
                                    "headers": table_data[0] if table_data else [],
                                    "rows": len(table_data),
                                    "cols": len(table_data[0]) if table_data else 0
                                },
                                bbox=corresponding_element.bbox if corresponding_element else None,
                                page_number=page_num,
                                confidence=0.8,  # pdfplumber는 confidence 제공하지 않음
                                metadata={
                                    "extraction_method": "pdfplumber"
                                }
                            )
                            tables.append(content)

        except Exception as e:
            logger.debug(f"pdfplumber 표 추출 실패 (페이지 {page_num}): {e}")

        return tables

    def _select_best_tables(self,
                           camelot_tables: List[ExtractedContent],
                           pdfplumber_tables: List[ExtractedContent],
                           original_elements: List[LayoutElement]) -> List[ExtractedContent]:
        """최고 품질의 표 선택"""
        best_tables = []

        # Camelot 결과 우선 (정확도 기반)
        high_quality_camelot = [t for t in camelot_tables if t.confidence > 0.8]
        if high_quality_camelot:
            best_tables.extend(high_quality_camelot)
            return best_tables

        # pdfplumber 결과 사용
        if pdfplumber_tables:
            best_tables.extend(pdfplumber_tables)
            return best_tables

        # 둘 다 없으면 Camelot 결과라도
        if camelot_tables:
            best_tables.extend(camelot_tables)

        return best_tables

    async def _extract_diagram_elements(self,
                                      diagram_elements: List[LayoutElement],
                                      document_path: Union[str, Path]) -> List[ExtractedContent]:
        """다이어그램 요소 추출"""
        extracted_diagrams = []

        try:
            document_path = Path(document_path)

            # PPTX인 경우 python-pptx로 처리
            if document_path.suffix.lower() == '.pptx':
                extracted_diagrams = await self._extract_diagrams_from_pptx(diagram_elements, document_path)

            # 다른 형식은 이미지 기반 처리
            else:
                # TODO: 이미지 기반 다이어그램 추출 (YOLO 등)
                pass

        except Exception as e:
            logger.error(f"다이어그램 추출 실패: {e}")

        return extracted_diagrams

    async def _extract_diagrams_from_pptx(self,
                                        diagram_elements: List[LayoutElement],
                                        pptx_path: Path) -> List[ExtractedContent]:
        """PPTX에서 다이어그램 추출"""
        extracted_diagrams = []

        if not PYTHON_PPTX_AVAILABLE:
            logger.warning("python-pptx를 사용할 수 없습니다.")
            return extracted_diagrams

        try:
            prs = Presentation(str(pptx_path))

            for slide_num, slide in enumerate(prs.slides):
                slide_diagrams = await self._extract_slide_diagrams(slide, slide_num + 1)
                extracted_diagrams.extend(slide_diagrams)

        except Exception as e:
            logger.error(f"PPTX 다이어그램 추출 실패: {e}")

        return extracted_diagrams

    async def _extract_slide_diagrams(self, slide, slide_num: int) -> List[ExtractedContent]:
        """슬라이드에서 다이어그램 추출"""
        diagrams = []

        try:
            # 도형과 커넥터 수집
            shapes = []
            connectors = []

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # 일반 도형
                    shape_data = {
                        "id": shape.shape_id,
                        "type": "shape",
                        "name": shape.name,
                        "text": shape.text if hasattr(shape, 'text') else "",
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                    shapes.append(shape_data)

                elif shape.shape_type == MSO_SHAPE_TYPE.CONNECTOR:
                    # 커넥터
                    connector_data = {
                        "id": shape.shape_id,
                        "type": "connector",
                        "name": shape.name,
                        "begin_x": shape.begin_x if hasattr(shape, 'begin_x') else 0,
                        "begin_y": shape.begin_y if hasattr(shape, 'begin_y') else 0,
                        "end_x": shape.end_x if hasattr(shape, 'end_x') else 0,
                        "end_y": shape.end_y if hasattr(shape, 'end_y') else 0
                    }
                    connectors.append(connector_data)

            # 그래프 구조 생성
            if shapes or connectors:
                graph_data = self._build_graph_from_shapes(shapes, connectors)

                content = ExtractedContent(
                    content_type="diagram",
                    data={
                        "graph": graph_data,
                        "shapes": shapes,
                        "connectors": connectors,
                        "slide_number": slide_num
                    },
                    page_number=slide_num,
                    confidence=0.9,
                    metadata={
                        "extraction_method": "python_pptx",
                        "slide_number": slide_num
                    }
                )
                diagrams.append(content)

        except Exception as e:
            logger.debug(f"슬라이드 {slide_num} 다이어그램 추출 실패: {e}")

        return diagrams

    def _build_graph_from_shapes(self, shapes: List[Dict], connectors: List[Dict]) -> Dict[str, Any]:
        """도형과 커넥터에서 그래프 구조 생성"""
        graph = {
            "nodes": [],
            "edges": []
        }

        # 노드 생성
        for shape in shapes:
            node = {
                "id": f"node_{shape['id']}",
                "label": shape['text'] or shape['name'],
                "type": "process",  # 기본값, 추후 형태 분석으로 개선 가능
                "position": {
                    "x": shape['left'] + shape['width'] // 2,
                    "y": shape['top'] + shape['height'] // 2
                },
                "size": {
                    "width": shape['width'],
                    "height": shape['height']
                }
            }
            graph["nodes"].append(node)

        # 엣지 생성 (커넥터 기반)
        for connector in connectors:
            # 시작점과 끝점에 가장 가까운 노드 찾기
            start_node = self._find_nearest_node(
                graph["nodes"], connector['begin_x'], connector['begin_y']
            )
            end_node = self._find_nearest_node(
                graph["nodes"], connector['end_x'], connector['end_y']
            )

            if start_node and end_node and start_node != end_node:
                edge = {
                    "id": f"edge_{connector['id']}",
                    "from": start_node['id'],
                    "to": end_node['id'],
                    "type": "directed"  # 기본값
                }
                graph["edges"].append(edge)

        return graph

    def _find_nearest_node(self, nodes: List[Dict], x: float, y: float, threshold: float = 100) -> Optional[Dict]:
        """가장 가까운 노드 찾기"""
        min_distance = float('inf')
        nearest_node = None

        for node in nodes:
            node_x = node['position']['x']
            node_y = node['position']['y']

            distance = ((x - node_x) ** 2 + (y - node_y) ** 2) ** 0.5

            if distance < min_distance and distance < threshold:
                min_distance = distance
                nearest_node = node

        return nearest_node

    async def _extract_image_elements(self,
                                    image_elements: List[LayoutElement],
                                    document_path: Union[str, Path]) -> List[ExtractedContent]:
        """이미지 요소 추출"""
        extracted_images = []

        try:
            document_path = Path(document_path)

            if document_path.suffix.lower() == '.pdf':
                extracted_images = await self._extract_images_from_pdf(image_elements, document_path)

        except Exception as e:
            logger.error(f"이미지 추출 실패: {e}")

        return extracted_images

    async def _extract_images_from_pdf(self,
                                     image_elements: List[LayoutElement],
                                     pdf_path: Path) -> List[ExtractedContent]:
        """PDF에서 이미지 추출"""
        extracted_images = []

        try:
            import fitz

            doc = fitz.open(str(pdf_path))

            for element in image_elements:
                page_num = element.page_number - 1  # 0-based
                if page_num < len(doc):
                    page = doc[page_num]
                    image_list = page.get_images()

                    # 해당 영역의 이미지 찾기
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            base_image = doc.extract_image(xref)

                            content = ExtractedContent(
                                content_type="image",
                                data={
                                    "image_data": base_image["image"],
                                    "format": base_image["ext"],
                                    "width": base_image["width"],
                                    "height": base_image["height"],
                                    "xref": xref
                                },
                                bbox=element.bbox,
                                page_number=element.page_number,
                                confidence=element.confidence,
                                metadata={
                                    "extraction_method": "pymupdf",
                                    "image_index": img_index
                                }
                            )
                            extracted_images.append(content)

                        except Exception as e:
                            logger.debug(f"이미지 추출 실패: {e}")

            doc.close()

        except Exception as e:
            logger.error(f"PDF 이미지 추출 실패: {e}")

        return extracted_images